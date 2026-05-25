#!/usr/bin/env python3
"""Generate UniMind PPTs — wow-moment focused, premium dark styling."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand ──
C = {
    'bg':       RGBColor(0x0A, 0x0A, 0x14),
    'bg2':      RGBColor(0x10, 0x10, 0x20),
    'surface':  RGBColor(0x16, 0x16, 0x2A),
    'card':     RGBColor(0x1A, 0x1A, 0x32),
    'border':   RGBColor(0x2A, 0x2A, 0x4A),
    'indigo':   RGBColor(0x63, 0x66, 0xF1),
    'indigo_l': RGBColor(0x81, 0x8C, 0xF8),
    'cyan':     RGBColor(0x22, 0xD3, 0xEE),
    'amber':    RGBColor(0xF5, 0x9E, 0x0B),
    'emerald':  RGBColor(0x10, 0xB9, 0x81),
    'rose':     RGBColor(0xF4, 0x3F, 0x5E),
    'white':    RGBColor(0xF0, 0xF0, 0xF5),
    'muted':    RGBColor(0xA0, 0xA0, 0xB8),
    'dim':      RGBColor(0x60, 0x60, 0x78),
    'black':    RGBColor(0x00, 0x00, 0x00),
}
LOGO = os.path.join(os.path.dirname(__file__), '..', '..', 'frontend', 'public', 'Unimind_logo.png')
W, H = Inches(13.333), Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def bg(slide, color=None):
    """Solid dark background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C['bg']


def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(1), radius=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill or C['card']
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    return s


def accent_bar(slide, l, t, w=Inches(1.2), h=Pt(4), color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color or C['indigo']
    s.line.fill.background()


def txt(slide, l, t, w, h, text, size=18, color=None, bold=False, align=PP_ALIGN.LEFT, font='Helvetica Neue'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color or C['white']
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def multi_txt(slide, l, t, w, h, lines, size=15, color=None, spacing=Pt(8)):
    """Multiple lines in one text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color or C['muted']
        r.font.name = 'Helvetica Neue'
    return tb


def bullets(slide, l, t, w, h, items, size=15, color=None, bullet_color=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    bc = bullet_color or C['indigo']
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = "●  "
        r1.font.size = Pt(size - 3)
        r1.font.color.rgb = bc
        r1.font.name = 'Helvetica Neue'
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color or C['muted']
        r2.font.name = 'Helvetica Neue'


def header(slide, label, title, subtitle=None):
    """Standard slide header with label + title + accent bar."""
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(0.7), Inches(0.35), height=Inches(0.4))
    txt(slide, Inches(0.7), Inches(1.0), Inches(10), Inches(0.35),
        label, size=11, color=C['indigo_l'], font='Menlo')
    txt(slide, Inches(0.7), Inches(1.4), Inches(11), Inches(0.7),
        title, size=32, color=C['white'], bold=True)
    accent_bar(slide, Inches(0.7), Inches(2.15))
    if subtitle:
        txt(slide, Inches(0.7), Inches(2.4), Inches(11), Inches(0.5),
            subtitle, size=16, color=C['muted'])


def page_num(slide, n, total):
    txt(slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.3),
        f"{n}/{total}", size=9, color=C['dim'], align=PP_ALIGN.RIGHT, font='Menlo')


def stat_block(slide, l, t, num, label, sub=None, num_color=None):
    rect(slide, l, t, Inches(3.5), Inches(1.6), fill=C['surface'], line=C['border'])
    txt(slide, l, t + Inches(0.15), Inches(3.5), Inches(0.7),
        num, size=38, color=num_color or C['cyan'], bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l, t + Inches(0.85), Inches(3.5), Inches(0.35),
        label, size=13, color=C['muted'], align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, l, t + Inches(1.15), Inches(3.5), Inches(0.3),
            sub, size=10, color=C['dim'], align=PP_ALIGN.CENTER, font='Menlo')


def card(slide, l, t, w, h, title, body_lines=None, icon=None, title_color=None, top_accent=None):
    rect(slide, l, t, w, h, fill=C['card'], line=C['border'])
    if top_accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = top_accent
        bar.line.fill.background()
    y = t + Inches(0.25)
    if icon:
        txt(slide, l + Inches(0.25), y, Inches(0.5), Inches(0.4), icon, size=22)
        l2 = l + Inches(0.7)
    else:
        l2 = l + Inches(0.25)
    txt(slide, l2, y, w - Inches(0.6), Inches(0.35),
        title, size=16, color=title_color or C['white'], bold=True)
    if body_lines:
        bullets(slide, l + Inches(0.25), y + Inches(0.45), w - Inches(0.5), h - Inches(0.8),
                body_lines, size=13)


def table(slide, l, t, w, rh, headers, rows, col_w=None):
    nr = len(rows) + 1
    nc = len(headers)
    ts = slide.shapes.add_table(nr, nc, l, t, w, rh * nr)
    tbl = ts.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = C['surface']
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = C['indigo_l']
            p.font.bold = True
            p.font.name = 'Menlo'
            p.alignment = PP_ALIGN.CENTER
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = C['card'] if ri % 2 == 0 else C['bg2']
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = C['muted']
                p.font.name = 'Helvetica Neue'
                p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════
#  CUSTOMER PPT — 10 slides, wow-loop focused
# ═══════════════════════════════════════════════════

# ═══════════════════════════════════════════════════
#  CUSTOMER PPT — 14 slides, wow-loop focused
# ═══════════════════════════════════════════════════

def gen_customer():
    prs = new_prs()
    T = 14

    # ── 1. Cover ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C['bg'])
    o = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    o.fill.solid(); o.fill.fore_color.rgb = C['indigo']; o.fill.fore_color.brightness = -0.85; o.line.fill.background()
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(5.7), Inches(1.2), height=Inches(1.2))
    txt(s, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.0),
        "你的老师出题，AI 帮忙审，学生练完自动记住。", size=38, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(4.0), Inches(9.3), Inches(0.5),
        "出题 → 练习 → 沉淀 → 更好的题 — 一个越用越聪明的闭环", size=18, color=C['cyan'], align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(5.0), Inches(9.3), Inches(0.4),
        "培训机构的 AI 运营系统  |  北京融知高科 · UniMind.ai", size=14, color=C['dim'], align=PP_ALIGN.CENTER)
    page_num(s, 1, T)

    # ── 2. Problem ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "THE PROBLEM", "你正在面对的三个问题")
    for i, (title, color, items) in enumerate([
        ("出题靠人工", C['rose'], ["教研老师 60% 时间在出题、组卷、写解析", "月薪 1.5-2 万，一年 18-24 万", "出一道题平均 30 分钟，质量参差不齐"]),
        ("续费率上不去", C['amber'], ["所有学生刷同一套题，无个性化", "遗忘曲线被忽视，复习效率低", "续费率卡在 60-70%"]),
        ("教学靠感觉", C['dim'], ["不知道谁在哪里卡住了", "效果无法量化，家长不信任", "好老师的经验无法沉淀复用"]),
    ]):
        l = Inches(0.7 + i * 4.1)
        card(s, l, Inches(2.8), Inches(3.7), Inches(3.5), title, items, title_color=color, top_accent=color)
    page_num(s, 2, T)

    # ── 3. THE FLYWHEEL (WOW) ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "THE FLYWHEEL", "一个越用越聪明的闭环",
           "这是 UniMind 的核心价值 — 不是三个独立工具，是一个自我强化的飞轮")
    for i, (title, sub, color, desc) in enumerate([
        ("① 教师出题", "AI 对抗出题引擎", C['indigo'],
         "输入考点 → 4 Agent 互相审核\n85%+ 可直接用于课堂\n一道题 10 秒，不是 30 分钟"),
        ("② 学生练习", "自适应刷题 + AI 批改", C['cyan'],
         "主观题自动评分 + 详细解析\n难度 5 级可选，教师锁定\n每次答题产生训练数据"),
        ("③ Memorix 沉淀", "记忆算法持续学习", C['emerald'],
         "每次答题更新个人记忆模型\nWeibull 遗忘曲线 + 在线学习\n复习越用越准，续费率 ↑"),
    ]):
        l = Inches(0.7 + i * 4.1)
        rect(s, l, Inches(2.8), Inches(3.7), Inches(3.0), fill=C['card'], line=color, lw=Pt(2))
        txt(s, l + Inches(0.3), Inches(2.95), Inches(3.1), Inches(0.4), title, size=20, color=color, bold=True)
        txt(s, l + Inches(0.3), Inches(3.4), Inches(3.1), Inches(0.3), sub, size=12, color=C['dim'], font='Menlo')
        txt(s, l + Inches(0.3), Inches(3.85), Inches(3.1), Inches(1.5), desc, size=14, color=C['muted'])
        if i < 2:
            txt(s, l + Inches(3.7), Inches(3.8), Inches(0.4), Inches(0.5), "▶", size=20, color=C['indigo_l'], align=PP_ALIGN.CENTER)
    rect(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.8), fill=C['surface'], line=C['indigo'])
    txt(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.6),
        "⟳ 飞轮效应：学生练习数据 → Memorix 发现薄弱点 → AI 生成更精准的题 → 效率更高 → 更多数据",
        size=14, color=C['cyan'], align=PP_ALIGN.CENTER)
    page_num(s, 3, T)

    # ── 4. ARC Pipeline ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "STEP ①", "教师出题：4-Agent 对抗引擎（ARC）",
           "不是「AI 随便出题」，是 4 个 Agent 互相审核，教师只管审核通过")
    for i, (label, color) in enumerate([("考点\n输入", C['dim']), ("Author\n出题", C['indigo']),
             ("Reviewer\n审核+工具", C['cyan']), ("修正", C['indigo']),
             ("Classifier\n分类审计", C['indigo']), ("入库\n可用", C['emerald'])]):
        l = Inches(0.4 + i * 2.1)
        rect(s, l, Inches(3.0), Inches(1.9), Inches(1.1), fill=color)
        txt(s, l, Inches(3.1), Inches(1.9), Inches(0.9), label, size=12, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
        if i < 5:
            txt(s, l + Inches(1.9), Inches(3.3), Inches(0.2), Inches(0.3), "→", size=13, color=C['dim'], align=PP_ALIGN.CENTER)
    txt(s, Inches(2.5), Inches(2.5), Inches(8), Inches(0.4),
        "↻ 最多 3 轮迭代 · 阈值 0.7 · 难度由教师锁定（AI 不能改）", size=12, color=C['dim'], align=PP_ALIGN.CENTER, font='Menlo')
    stat_block(s, Inches(0.7), Inches(4.7), "~60%", "单次 LLM 可用率", num_color=C['dim'])
    stat_block(s, Inches(4.9), Inches(4.7), "85%+", "ARC 管线可用率", num_color=C['emerald'])
    stat_block(s, Inches(9.1), Inches(4.7), "50x", "出题速度提升", "30min → 10s", C['cyan'])
    bullets(s, Inches(0.7), Inches(6.5), Inches(11.5), Inches(0.8),
            ["支持题型：单选 / 多选 / 判断 / 名词解释 / 简答 / 论述 / 计算",
             "5 级难度：入门 / 简单 / 一般 / 困难 / 地狱 — 教师手动锁定"], size=13)
    page_num(s, 4, T)

    # ── 5. Student Practice ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "STEP ②", "学生练习：扫码即练，AI 批改",
           "无需下载 APP，纯网页端，手机/平板/电脑全支持")
    for i, (title, sub) in enumerate([("生成邀请", "链接/二维码"), ("扫码注册", "填写信息"), ("开始刷题", "AI 自适应"), ("AI 批改", "主观题也行")]):
        l = Inches(0.7 + i * 3.1)
        rect(s, l, Inches(3.0), Inches(2.7), Inches(1.1), fill=C['indigo'] if i >= 2 else C['surface'])
        txt(s, l, Inches(3.05), Inches(2.7), Inches(0.5), title, size=16, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
        txt(s, l, Inches(3.55), Inches(2.7), Inches(0.4), sub, size=12, color=C['indigo_l'], align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, l + Inches(2.7), Inches(3.3), Inches(0.4), Inches(0.3), "→", size=16, color=C['dim'], align=PP_ALIGN.CENTER)
    bullets(s, Inches(0.7), Inches(4.7), Inches(5.5), Inches(2.5), [
        "主观题 AI 评分 + 详细解析，教师可复核",
        "5 级难度：入门/简单/一般/困难/地狱",
        "教师锁定难度，AI 不能擅自更改",
        "ELO 排位系统，激发学习动力",
        "每周自动生成「认知资产周报」",
        "中英文双语界面",
    ], size=14)
    card(s, Inches(7.0), Inches(4.7), Inches(5.5), Inches(2.5), "学生端亮点", [
        "零门槛：扫码即练，无需下载",
        "即时反馈：答题后立即看到解析",
        "个性化：系统记住每个学生的薄弱点",
        "社交：ELO 排名 + 自习室互动",
    ], top_accent=C['cyan'])
    page_num(s, 5, T)

    # ── 6. Memorix ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "STEP ③", "Memorix 沉淀：越用越准的记忆算法",
           "不是固定间隔复习，是为每个学生实时学习的自研算法")
    table(s, Inches(0.7), Inches(2.8), Inches(11.5), Inches(0.42),
          ["", "传统间隔重复", "FSRS v4.5（开源最优）", "Memorix（UniMind）"],
          [
              ["遗忘模型", "固定曲线", "幂律分布", "Weibull（可学习）"],
              ["个性化", "无", "全局参数", "每用户 20 维权重"],
              ["更新方式", "无", "批量拟合", "每次答题实时学习"],
              ["预测精度", "差", "RMSE 基线", "↓ 13.7%"],
              ["用户留存", "基线", "基线", "↑ 9.2%"],
          ],
          col_w=[Inches(2.0), Inches(2.5), Inches(3.2), Inches(3.8)])
    rect(s, Inches(0.7), Inches(5.4), Inches(11.5), Inches(1.3), fill=C['surface'], line=C['emerald'])
    txt(s, Inches(1.0), Inches(5.5), Inches(11.0), Inches(0.4),
        "⟳ 飞轮闭环", size=16, color=C['emerald'], bold=True)
    txt(s, Inches(1.0), Inches(5.9), Inches(11.0), Inches(0.7),
        "学生每做一道题 → Memorix 更新个人模型 → 下次复习更精准 → 薄弱点自动暴露 → AI 生成更有针对性的题 → 效率更高 → 续费率提升",
        size=14, color=C['cyan'])
    txt(s, Inches(0.7), Inches(6.9), Inches(11), Inches(0.3),
        "验证：431 金融考研 · 500+ 用户 · 120,000+ 复习记录 · 白皮书已发布", size=11, color=C['dim'], font='Menlo')
    page_num(s, 6, T)

    # ── 7. Knowledge Graph ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "FEATURE", "知识图谱 + 学情分析",
           "红黄绿一眼看清，谁在哪个知识点上卡住了")
    bullets(s, Inches(0.7), Inches(2.8), Inches(5.5), Inches(2.5), [
        "4 层知识树：学科 → 章 → 节 → 知识点",
        "每个知识点：掌握度评分 + 热力图（绿/黄/红）",
        "自动识别瓶颈：DAG 遍历找到前置知识缺口",
        "用户自标注：优先级、置信度、标签、笔记",
        "支持 Markdown 导入 + AI 自动生成知识树",
    ], size=14)
    card(s, Inches(7.0), Inches(2.8), Inches(5.5), Inches(2.5), "场景举例", [
        "学生在「投资组合管理」反复出错",
        "系统发现「资产定价模型」掌握度仅 35%",
        "自动推荐回到 CAPM 复习",
        "教师在后台看到全班瓶颈分布",
    ], top_accent=C['emerald'])
    txt(s, Inches(0.7), Inches(5.8), Inches(11.5), Inches(0.5),
        "知识图谱是飞轮的数据骨架 — 每个知识点的记忆状态由 Memorix 驱动，教师出题时可按图谱选考点",
        size=14, color=C['cyan'])
    page_num(s, 7, T)

    # ── 8. AI Assistant + Interview ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "FEATURE", "AI 助教 + 模拟面试")
    card(s, Inches(0.7), Inches(2.8), Inches(5.8), Inches(4.0), "AI 助教", [
        "多角色可选：不同性格和学术背景",
        "支持 LaTeX 公式渲染（数学/金融/物理）",
        "7×24 在线，学生随时提问",
        "自动查询知识树 + 学生薄弱点 + 错题",
        "机构可自定义 Bot 人设和知识库",
        "WebSocket 流式响应，打字机效果",
    ], top_accent=C['indigo'])
    card(s, Inches(7.0), Inches(2.8), Inches(5.5), Inches(4.0), "AI 模拟面试", [
        "4 种面试类型：简历/英语/专业/综合",
        "2 种风格：友好引导 / 压力面试",
        "简历上传 → AI 优化 + 预测 10 个提问",
        "实时 WebSocket 对话，逐轮反馈",
        "面试结束 → 5 维雷达图评分",
        "维度：理论深度/逻辑/抗压/表达/英语",
    ], top_accent=C['cyan'])
    page_num(s, 8, T)

    # ── 9. More Features ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "FEATURES", "更多功能")
    for i, (title, desc) in enumerate([
        ("视频课程", "分片上传 + ASR 转录 + AI 大纲\n时间戳跳转 + 字幕叠加 + 倍速"),
        ("在线自习室", "番茄钟 + 实时聊天 + 周计划\n学习状态广播 + GIPHY + 图片"),
        ("答疑系统", "社区式问答，教师认证标识\n支持图片、搜索、点赞"),
        ("PDF 模考", "AI 根据薄弱点自动组卷\n生成 PDF + 答案，支持教师发布"),
        ("学情周报", "ELO 百分位 + 永久掌握题数\n准确率/题量/专注时长趋势图"),
        ("通知系统", "站内消息 + 管理员批量推送\nWebSocket 实时送达"),
    ]):
        col, row = i % 3, i // 3
        l = Inches(0.7 + col * 4.1)
        t = Inches(2.8 + row * 2.2)
        card(s, l, t, Inches(3.7), Inches(1.9), title, [desc])
    page_num(s, 9, T)

    # ── 10. Admin Panel ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "ADMIN", "机构管理后台", "一个后台管所有事，教师无需技术背景")
    for i, (title, desc) in enumerate([
        ("题库管理", "批量导入 / AI 生成 / 审核队列 / 标签管理 / 难度标记"),
        ("学员管理", "分班 / 数据导出 / 学情对比 / ELO 排名 / 激活码"),
        ("课程管理", "视频上传 / AI 大纲 / 课程标签 / 分片上传 / 进度跟踪"),
        ("AI Bot 管理", "自定义助教人设 / 知识库配置 / 多角色切换"),
        ("Prompt 管理", "文件级模板 / 数据库版本历史 / 一键回滚 / 在线编辑"),
        ("数据洞察", "班级对比 / 知识点掌握分布 / 使用量统计 / 审计日志"),
    ]):
        col, row = i % 2, i // 2
        l = Inches(0.7 + col * 6.2)
        t = Inches(2.8 + row * 1.5)
        txt(s, l, t, Inches(2.5), Inches(0.35), title, size=16, color=C['indigo_l'], bold=True)
        txt(s, l + Inches(0.3), t + Inches(0.4), Inches(5.5), Inches(0.4), desc, size=13, color=C['muted'])
    txt(s, Inches(0.7), Inches(7.0), Inches(11), Inches(0.3),
        "三层权限：平台角色 ∩ 机构角色 ∩ 方案等级 = 有效权限 · 多租户数据隔离", size=11, color=C['dim'], font='Menlo')
    page_num(s, 10, T)

    # ── 11. Subject Coverage ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "COVERAGE", "学科覆盖", "学科无关 — 只要能结构化为知识点，就能用")
    subjects = [
        ("考研专业课", "金融 431 / 法学 / 医学 / 408 / 311 / 312"),
        ("职业资格证", "CPA / CFA / 法考 / 教资 / 一建 / USMLE"),
        ("中学学科", "高中数学 / 物理 / 化学 / 生物"),
        ("公考", "行测 / 申论 / 公基 / 军队文职"),
        ("留学考试", "SAT / ACT / AP / GRE / GMAT / LSAT"),
        ("语言考试", "TOEFL / IELTS"),
    ]
    for i, (cat, items) in enumerate(subjects):
        col, row = i % 3, i // 3
        l = Inches(0.7 + col * 4.1)
        t = Inches(2.8 + row * 2.2)
        rect(s, l, t, Inches(3.7), Inches(1.9), fill=C['card'], line=C['border'])
        txt(s, l + Inches(0.3), t + Inches(0.2), Inches(3.1), Inches(0.35), cat, size=16, color=C['indigo_l'], bold=True)
        txt(s, l + Inches(0.3), t + Inches(0.65), Inches(3.1), Inches(1.0), items, size=14, color=C['muted'])
    page_num(s, 11, T)

    # ── 12. Pricing ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "PRICING", "定价方案")
    table(s, Inches(0.4), Inches(2.8), Inches(12.5), Inches(0.4),
          ["", "个人工作室", "独立教师", "培训机构", "连锁品牌"],
          [
              ["价格", "免费", "¥499/月", "¥1,299/月", "¥3,999/月"],
              ["学生数", "30", "50", "200", "不限"],
              ["教师数", "1", "1", "5", "不限"],
              ["AI 出题", "30 次/月", "100 次/月", "不限", "不限"],
              ["AI 批改", "不限", "不限", "不限", "不限"],
              ["Memorix", "—", "✓", "✓", "✓"],
              ["知识图谱", "—", "—", "✓", "✓"],
              ["视频课程", "—", "—", "✓", "✓"],
              ["品牌定制", "—", "—", "—", "✓"],
              ["私有部署", "—", "—", "—", "✓"],
          ],
          col_w=[Inches(1.8), Inches(2.4), Inches(2.6), Inches(2.9), Inches(2.8)])
    txt(s, Inches(0.7), Inches(6.9), Inches(11), Inches(0.3),
        "年付约 8 折（节省 23-33%）· 14 天免费试用 · 无需绑卡 · 数据永久保留", size=12, color=C['dim'])
    page_num(s, 12, T)

    # ── 13. Testimonials + Stats ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "PROOF", "已在使用的机构")
    for i, (name, quote) in enumerate([
        ("何伟 / 北京 / CFA 培训", "生成了 400+ 道题，85% 以上可以直接用于课堂教学。"),
        ("周鹏 / 成都 / 高中数学连锁", "续费率从不到 70% 提升到 80% 以上。3 校区、200+ 学生。"),
        ("林佳 / 上海 / CPA 培训", "当天注册当天就能用，不需要技术团队。"),
        ("吴敏 / 广州 / 法考培训", "批改时间缩短了一半以上。"),
    ]):
        col, row = i % 2, i // 2
        l = Inches(0.7 + col * 6.2)
        t = Inches(2.8 + row * 2.0)
        rect(s, l, t, Inches(5.8), Inches(1.7), fill=C['card'], line=C['border'])
        txt(s, l + Inches(0.3), t + Inches(0.2), Inches(5.2), Inches(0.8), f'“{quote}”', size=14, color=C['muted'])
        txt(s, l + Inches(0.3), t + Inches(1.1), Inches(5.2), Inches(0.3), f"— {name}", size=12, color=C['indigo_l'], bold=True)
    stat_block(s, Inches(0.7), Inches(7.0), "", "", "")
    txt(s, Inches(2.5), Inches(7.0), Inches(8), Inches(0.3),
        "10+ 学科  ·  50,000+ 已生成题目  ·  50+ 活跃机构  ·  50x 出题提速", size=14, color=C['dim'], align=PP_ALIGN.CENTER)
    page_num(s, 13, T)

    # ── 14. CTA ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C['bg'])
    o = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    o.fill.solid(); o.fill.fore_color.rgb = C['indigo']; o.fill.fore_color.brightness = -0.85; o.line.fill.background()
    txt(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.8),
        "让你的老师专注教学，把运营交给 AI", size=36, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(3.0), Inches(9.3), Inches(0.5),
        "一个越用越聪明的闭环：出题 → 练习 → 沉淀 → 更好的题", size=18, color=C['cyan'], align=PP_ALIGN.CENTER)
    for i, line in enumerate([
        "① 发邮件至 korsonedu@gmail.com 申请试用",
        "② 14 天 Growth 方案免费体验（无需绑卡）",
        "③ 专属对接人全程陪跑，当天注册当天能用",
    ]):
        txt(s, Inches(3.5), Inches(4.2 + i * 0.55), Inches(6.3), Inches(0.4), line, size=18, color=C['indigo_l'], align=PP_ALIGN.CENTER)
    rect(s, Inches(4.2), Inches(5.8), Inches(4.9), Inches(0.6), fill=C['amber'])
    txt(s, Inches(4.2), Inches(5.85), Inches(4.9), Inches(0.5),
        "⭐ 首批机构免费用 Growth 方案 · 截止 2026 年 6 月 30 日", size=14, color=C['bg'], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(6.8), Inches(9.3), Inches(0.3),
        "北京融知高科 · UniMind.ai · 培训机构的 AI 基础设施", size=11, color=C['dim'], align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(7.1), Inches(9.3), Inches(0.3),
        "前 20 名付费客户锁定终身价格", size=12, color=C['amber'], align=PP_ALIGN.CENTER)
    page_num(s, 14, T)

    out = os.path.join(os.path.dirname(__file__), 'UniMind_客户版.pptx')
    prs.save(out)
    print(f"Saved: {out}")


# ═══════════════════════════════════════════════════
#  INVESTOR PPT — 16 slides, wow-loop + moat focused
# ═══════════════════════════════════════════════════

def gen_investor():
    prs = new_prs()
    T = 16

    # ── 1. Cover ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C['bg'])
    o = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    o.fill.solid(); o.fill.fore_color.rgb = C['indigo']; o.fill.fore_color.brightness = -0.85; o.line.fill.background()
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(5.7), Inches(0.8), height=Inches(1.2))
    txt(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.8), "UniMind.ai", size=52, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(3.5), Inches(9.3), Inches(0.5), "培训机构的 AI 运营系统", size=22, color=C['cyan'], align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(4.2), Inches(9.3), Inches(0.4), "The Operations Layer for Tutoring Businesses", size=14, color=C['dim'], align=PP_ALIGN.CENTER, font='Menlo')
    txt(s, Inches(2), Inches(5.2), Inches(9.3), Inches(0.4), "教师出题 → 学生练习 → Memorix 沉淀 → 越用越聪明的飞轮", size=16, color=C['indigo_l'], align=PP_ALIGN.CENTER)
    txt(s, Inches(2), Inches(6.2), Inches(9.3), Inches(0.3), "北京融知高科", size=13, color=C['dim'], align=PP_ALIGN.CENTER)
    page_num(s, 1, T)

    # ── 2. Market ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "MARKET", "万亿教培市场的效率黑洞", "中国教培市场 ~8,000 亿（2025），中小机构占 60%+，无技术团队")
    for i, (num, label, sub, color) in enumerate([
        ("¥2 万/月", "教研人力成本", "60% 时间在出题", C['rose']),
        ("60-70%", "行业平均续费率", "缺少个性化手段", C['amber']),
        ("0", "实时学情数据", "决策靠感觉", C['dim']),
    ]):
        stat_block(s, Inches(0.7 + i * 4.1), Inches(3.2), num, label, sub, color)
    bullets(s, Inches(0.7), Inches(5.4), Inches(11.5), Inches(1.5), [
        "中小机构（10-50 人）占市场 60%+，高度分散，绝大多数无技术团队",
        "考研/公考/K12/职业资格 — 每个赛道都有数万家小机构",
        "大厂 SaaS（校宝、小鹅通）解决管理问题，不解决教学效率问题",
        "AI 教育工具（松鼠 AI 等）面向 C 端，不面向 B 端机构",
    ], size=14)
    page_num(s, 2, T)

    # ── 3. Flywheel (WOW) ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "THE FLYWHEEL", "核心价值：越用越聪明的闭环", "不是三个独立工具，是一个自我强化的数据飞轮")
    for i, (title, sub, color, desc) in enumerate([
        ("① 教师 AI 出题", "ARC 4-Agent 对抗引擎", C['indigo'],
         "输入考点 → 4 Agent 互相审核\n85%+ 可直接用于课堂\n出题速度 50x 提升"),
        ("② 学生自适应练习", "AI 批改 + ELO 排位", C['cyan'],
         "主观题自动评分\n难度教师锁定\n每次答题产生训练数据"),
        ("③ Memorix 记忆沉淀", "Weibull + 在线 SGD", C['emerald'],
         "20 维个人权重实时更新\n预测精度比 FSRS ↑13.7%\n用户留存 ↑9.2%"),
    ]):
        l = Inches(0.7 + i * 4.1)
        rect(s, l, Inches(2.8), Inches(3.7), Inches(3.0), fill=C['card'], line=color, lw=Pt(2))
        txt(s, l + Inches(0.3), Inches(2.95), Inches(3.1), Inches(0.35), title, size=18, color=color, bold=True)
        txt(s, l + Inches(0.3), Inches(3.35), Inches(3.1), Inches(0.25), sub, size=11, color=C['dim'], font='Menlo')
        txt(s, l + Inches(0.3), Inches(3.75), Inches(3.1), Inches(1.4), desc, size=14, color=C['muted'])
        if i < 2:
            txt(s, l + Inches(3.7), Inches(3.7), Inches(0.4), Inches(0.4), "→", size=20, color=C['indigo_l'], align=PP_ALIGN.CENTER)
    rect(s, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.9), fill=C['surface'], line=C['indigo'])
    txt(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.35),
        "⟳ 飞轮：学生数据 → Memorix 发现薄弱点 → AI 生成更精准的题 → 效率 ↑ → 续费 ↑ → 更多数据",
        size=14, color=C['cyan'], align=PP_ALIGN.CENTER)
    txt(s, Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.3),
        "数据飞轮 = 护城河。竞品只做了其中一环，UniMind 打通全链路。", size=13, color=C['amber'], align=PP_ALIGN.CENTER)
    page_num(s, 3, T)

    # ── 4. Memorix Moat ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "MOAT 01", "Memorix 自适应记忆算法", "自研间隔重复算法 · 白皮书级严谨度 · 6 大创新")
    for i, (title, desc) in enumerate([
        ("Weibull 遗忘模型", "R(t)=exp(-(t/λ)^k)\nk 可学习，比幂律更贴合"),
        ("在线 SGD + 力量", "每次答题后实时更新\nNesterov 力量加速收敛"),
        ("贝叶斯先验", "20 维个性化权重\nL2 正则 → 冷启动快"),
        ("Brier Score 损失", "严格 proper scoring\n概率校准保证"),
        ("遗憾最小化调度", "平衡遗忘风险\n与复习机会成本"),
        ("知识嵌入", "8 维向量 + 用户对齐\n跨知识点迁移"),
    ]):
        col, row = i % 3, i // 3
        l = Inches(0.7 + col * 4.1)
        t = Inches(2.8 + row * 2.1)
        card(s, l, t, Inches(3.7), Inches(1.8), title, [desc])
    txt(s, Inches(0.7), Inches(7.0), Inches(11), Inches(0.3),
        "验证：431 金融考研 · 500+ 用户 · 120K+ 记录 · RMSE -13.7% · 留存 +9.2%", size=11, color=C['dim'], font='Menlo')
    page_num(s, 4, T)

    # ── 5. ARC Moat ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "MOAT 02", "4-Agent 对抗出题管线（ARC）", "Reviewer 带自主研究工具，可查询知识库后审核")
    for i, (label, color) in enumerate([("Author\n生成", C['indigo']), ("Reviewer\n审核+工具", C['cyan']),
             ("AuthorRevise\n修正", C['indigo']), ("Classifier\n审计", C['indigo'])]):
        l = Inches(0.7 + i * 3.1)
        rect(s, l, Inches(3.0), Inches(2.7), Inches(1.2), fill=color)
        txt(s, l, Inches(3.1), Inches(2.7), Inches(1.0), label, size=15, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, l + Inches(2.7), Inches(3.3), Inches(0.4), Inches(0.3), "→", size=16, color=C['dim'], align=PP_ALIGN.CENTER)
    txt(s, Inches(3), Inches(4.4), Inches(7), Inches(0.3),
        "最多 3 轮 · 阈值 0.7 · 难度教师锁定（AI 不能改）", size=12, color=C['dim'], align=PP_ALIGN.CENTER, font='Menlo')
    bullets(s, Inches(0.7), Inches(5.0), Inches(11.5), Inches(1.8), [
        "Reviewer 可自主调用 lookup_definition() 和 search_similar() 查询知识库",
        "3 维评分：区分度、清晰度、知识覆盖度",
        "难度控制外生：教师锁定 → Classifier 检测偏差 → 人工审核队列",
    ], size=14)
    page_num(s, 5, T)

    # ── 6. AI Engine ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "MOAT 03", "AI 引擎架构", "Provider 无关，热切换，只改 2 个常量")
    for i, (title, desc) in enumerate([
        ("Prompt as Code", "文件模板 + DB 版本历史 + 一键回滚"),
        ("Schema 强制输出", "tool_choice=\"required\" + JSON Schema，零解析失败"),
        ("模型路由集中化", "config.py 单点控制，provider 热切换"),
        ("熔断器模式", "5 次失败 → 熔断 → 5min 后半开探测"),
        ("降级策略", "AI 失败 → 人工题库 / 标准答案 / 保存进度"),
        ("模型分层", "fast（对话/出题）+ pro（审核/判分/知识树）"),
    ]):
        t = Inches(2.8 + i * 0.72)
        txt(s, Inches(0.7), t, Inches(3.0), Inches(0.35), f"●  {title}", size=14, color=C['indigo_l'], bold=True)
        txt(s, Inches(4.0), t, Inches(8.5), Inches(0.35), desc, size=13, color=C['muted'])
    txt(s, Inches(0.7), Inches(7.0), Inches(11), Inches(0.3),
        "Provider 无关：当前 MiMo V2.5，切换只需改 2 个常量 · thinking + tool_choice 不冲突", size=11, color=C['dim'], font='Menlo')
    page_num(s, 6, T)

    # ── 7. Product Overview ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "PRODUCT", "产品全景")
    for i, (text, fill, color) in enumerate([
        ("学员端 · 教师端 · 管理端 · 机构品牌页", C['surface'], C['muted']),
        ("AI 出题 · 自适应复习 · AI 批改 · 知识图谱 · AI 助教 · 视频 · 面试 · 自习室 · 答疑 · PDF 模考 · 学情报告", C['indigo'], C['white']),
        ("AI 引擎 · 熔断器 · WebSocket · RBAC · 加密字段 · 多租户", C['bg2'], C['muted']),
        ("Django 6.0 + React 19 + PostgreSQL | Celery + Redis + Daphne", C['card'], C['dim']),
    ]):
        t = Inches(2.8 + i * 1.05)
        rect(s, Inches(0.7), t, Inches(11.5), Inches(0.85), fill=fill)
        txt(s, Inches(0.9), t + Inches(0.15), Inches(11.0), Inches(0.55), text, size=13, color=color, align=PP_ALIGN.CENTER)
    page_num(s, 7, T)

    # ── 8. Tech Stack ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "TECH", "技术栈")
    for i, (label, value) in enumerate([
        ("前端", "React 19 + TypeScript + Vite + shadcn/ui + Tailwind 4"),
        ("后端", "Django 6.0 + DRF 3.16 + Python 3.12+"),
        ("实时通信", "Django Channels 4.3 + Daphne 4.2（WebSocket）"),
        ("异步任务", "Celery 5.4 + Redis"),
        ("数据库", "PostgreSQL（生产）/ SQLite（开发）"),
        ("AI 引擎", "MiMo V2.5（可热切换任意 LLM provider）"),
        ("部署", "Nginx + Daphne x3 + Celery Worker x2 + systemd"),
        ("CI/CD", "GitHub Actions（pip-audit + npm audit + build + smoke test）"),
        ("i18n", "i18next 中英文双语（20 namespace）"),
    ]):
        t = Inches(2.6 + i * 0.5)
        txt(s, Inches(0.7), t, Inches(2.0), Inches(0.35), label, size=13, color=C['indigo_l'], bold=True)
        txt(s, Inches(3.0), t, Inches(9.5), Inches(0.35), value, size=13, color=C['muted'])
    page_num(s, 8, T)

    # ── 9. Security ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "SECURITY", "安全与合规")
    for i, (title, desc) in enumerate([
        ("认证", "Cookie-first Token 认证 + Authorization header fallback"),
        ("加密", "Fernet AES 加密支付密钥、敏感字段，ENCRYPTION_KEY 环境变量"),
        ("权限", "三层 RBAC：平台角色 ∩ 机构角色 ∩ 方案等级 = 有效权限"),
        ("多租户隔离", "queryset 级别数据隔离，机构用户只能看自己数据"),
        ("安全头", "HSTS / X-Frame-Options DENY / X-Content-Type-Options nosniff"),
        ("熔断器", "AI 服务连续失败自动熔断，降级到人工题库"),
        ("限流", "登录/注册/验证码接口滑动窗口限流"),
        ("备份", "自动 pg_dump 压缩备份，30 天自动清理"),
    ]):
        col, row = i % 2, i // 2
        l = Inches(0.7 + col * 6.2)
        t = Inches(2.6 + row * 1.15)
        txt(s, l, t, Inches(2.0), Inches(0.35), f"●  {title}", size=14, color=C['indigo_l'], bold=True)
        txt(s, l + Inches(0.3), t + Inches(0.35), Inches(5.5), Inches(0.35), desc, size=12, color=C['muted'])
    page_num(s, 9, T)

    # ── 10. Business Model ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "BUSINESS", "商业模式", "SaaS 订阅 + 4 级定价")
    table(s, Inches(0.4), Inches(2.8), Inches(12.5), Inches(0.42),
          ["方案", "月价", "学生", "教师", "目标", "年化 ARPU"],
          [
              ["Free", "¥0", "30", "1", "工作室", "¥0"],
              ["Starter", "¥499", "50", "1", "独立教师", "¥5,988"],
              ["Growth", "¥1,299", "200", "5", "培训机构", "¥15,588"],
              ["Enterprise", "¥3,999", "不限", "不限", "连锁品牌", "¥47,988"],
          ],
          col_w=[Inches(1.6), Inches(1.3), Inches(1.2), Inches(1.2), Inches(2.5), Inches(2.0)])
    card(s, Inches(0.7), Inches(5.2), Inches(5.8), Inches(1.8), "单位经济学（Growth）",
         ["LTV（24 月）：¥31,176", "CAC：~¥2,000-5,000", "LTV/CAC：6-19x（健康）"])
    card(s, Inches(7.0), Inches(5.2), Inches(5.8), Inches(1.8), "增长引擎",
         ["Freemium 漏斗：Free → Starter → Growth", "学生数 ↑ → 方案升级", "年付 8 折锁定现金流"])
    page_num(s, 10, T)

    # ── 11. GTM ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "GTM", "市场策略", "分阶段 GTM")
    for i, (phase, target, strategy, color) in enumerate([
        ("P0 · 0-6 月", "10+ 付费 · ¥50 万 ARR", "考研/公考中小机构\n14 天试用 + 人工陪跑", C['emerald']),
        ("P1 · 6-12 月", "50+ 机构 · ¥200 万 ARR", "K12 + 职业资格\n内容营销 + 渠道合作", C['indigo']),
        ("P2 · 12-24 月", "200+ 机构 · ¥800 万 ARR", "连锁品牌 + 出海\nStripe USD 结算", C['cyan']),
    ]):
        l = Inches(0.7 + i * 4.1)
        rect(s, l, Inches(2.8), Inches(3.7), Inches(2.3), fill=C['card'], line=color, lw=Pt(2))
        txt(s, l + Inches(0.3), Inches(2.95), Inches(3.1), Inches(0.3), phase, size=14, color=color, bold=True)
        txt(s, l + Inches(0.3), Inches(3.3), Inches(3.1), Inches(0.3), target, size=16, color=C['white'], bold=True)
        txt(s, l + Inches(0.3), Inches(3.7), Inches(3.1), Inches(0.9), strategy, size=13, color=C['muted'])
    for i, step in enumerate(["线索 2-3/周", "试用注册", "学员激活 >50%", "付费 >20%", "续费 >90%/月"]):
        l = Inches(0.5 + i * 2.5)
        rect(s, l, Inches(5.5), Inches(2.2), Inches(0.65), fill=C['indigo'] if i > 0 else C['surface'])
        txt(s, l, Inches(5.55), Inches(2.2), Inches(0.55), step, size=11, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
        if i < 4:
            txt(s, l + Inches(2.2), Inches(5.6), Inches(0.3), Inches(0.3), "→", size=12, color=C['dim'], align=PP_ALIGN.CENTER)
    bullets(s, Inches(0.7), Inches(6.5), Inches(11.5), Inches(0.8), [
        "获客渠道：教培社群内容营销 · 首批客户口碑 · 机构首页 SEO · 限时免费政策",
    ], size=13)
    page_num(s, 11, T)

    # ── 12. Competition ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "COMPETITION", "竞争格局")
    table(s, Inches(0.4), Inches(2.8), Inches(12.5), Inches(0.42),
          ["", "教培 SaaS（校宝）", "AI 教育（松鼠 AI）", "UniMind"],
          [
              ["定位", "管理工具", "学习工具", "AI 运营系统"],
              ["AI 能力", "无/弱", "单一", "多 Agent 管线"],
              ["出题", "无", "质量不可控", "ARC 85%+"],
              ["个性化", "无", "有", "Memorix 在线学习"],
              ["B2B 管理", "强", "弱", "强（RBAC/白标）"],
              ["自研算法", "无", "有", "有（白皮书+工程）"],
          ],
          col_w=[Inches(2.0), Inches(3.0), Inches(3.0), Inches(4.5)])
    txt(s, Inches(0.7), Inches(5.8), Inches(11.5), Inches(0.4),
        "护城河：Memorix 算法 · ARC 管线 · 全链路整合 · 数据飞轮", size=14, color=C['amber'], bold=True, align=PP_ALIGN.CENTER)
    page_num(s, 12, T)

    # ── 13. Traction ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "TRACTION", "牵引力与里程碑")
    card(s, Inches(0.7), Inches(2.8), Inches(5.8), Inches(2.2), "已达成", [
        "产品已上线运营（UniMind.ai，阿里云上海）",
        "10+ 学科支持（考研、公考、K12、职业资格）",
        "50,000+ 题目已生成 · 50+ 活跃机构",
        "完整 B2B SaaS + 自研算法 + 白皮书",
    ], top_accent=C['emerald'])
    card(s, Inches(7.0), Inches(2.8), Inches(5.8), Inches(2.2), "近期目标（6 月）", [
        "10+ 付费机构客户",
        "¥50 万 ARR",
        "首批客户案例 + 数据验证",
    ], top_accent=C['indigo'])
    card(s, Inches(0.7), Inches(5.4), Inches(5.8), Inches(1.5), "中期目标（12 月）", [
        "50+ 付费机构 · ¥200 万 ARR",
        "出海准备（Stripe + 英文市场验证）",
    ], top_accent=C['cyan'])
    card(s, Inches(7.0), Inches(5.4), Inches(5.8), Inches(1.5), "关键指标", [
        "试用转化 >20% · 月留存 >90% · 学员激活 >50%",
    ], top_accent=C['amber'])
    page_num(s, 13, T)

    # ── 14. Financial ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "FINANCIALS", "财务预测（保守估计）")
    table(s, Inches(1.5), Inches(2.8), Inches(10.0), Inches(0.45),
          ["指标", "6 个月", "12 个月", "24 个月"],
          [
              ["付费机构", "10", "50", "200"],
              ["平均 ARPU", "¥15,000/年", "¥18,000/年", "¥20,000/年"],
              ["ARR", "¥15 万", "¥90 万", "¥400 万"],
              ["月度留存", "90%", "92%", "95%"],
              ["毛利率", "70%", "75%", "80%"],
          ],
          col_w=[Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)])
    card(s, Inches(1.5), Inches(5.5), Inches(10.0), Inches(1.5), "关键假设", [
        "LLM 成本随规模下降（当前 ~¥0.01/次 API 调用）",
        "基础设施成本固定（单服务器可支撑 100+ 机构）",
        "人力成本主要是产品迭代和客户成功",
    ])
    page_num(s, 14, T)

    # ── 15. Team + Funding ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    header(s, "TEAM & ASK", "团队与融资")
    card(s, Inches(0.7), Inches(2.8), Inches(5.8), Inches(4.0), "团队", [
        "— 创始人背景（教育 + 技术交叉）",
        "— 核心技术能力（全栈 + AI 算法）",
        "— 行业认知（对教培痛点的理解）",
    ], top_accent=C['indigo'])
    card(s, Inches(7.0), Inches(2.8), Inches(5.8), Inches(4.0), "融资需求", [
        "— 本轮融资金额和轮次",
        "— 资金用途：",
        "    40% 产品迭代（出海/移动端）",
        "    30% 市场获客（种子客户 + 内容营销）",
        "    20% 团队扩充（客户成功 + 销售）",
        "    10% 运营储备",
        "— 预期达成里程碑",
    ], top_accent=C['cyan'])
    page_num(s, 15, T)

    # ── 16. Vision ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C['bg'])
    o = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
    o.fill.solid(); o.fill.fore_color.rgb = C['indigo']; o.fill.fore_color.brightness = -0.85; o.line.fill.background()
    txt(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.8),
        "让每一个培训机构都拥有 AI 运营能力", size=36, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
    for i, (stage, desc) in enumerate([
        ("今天", "AI 出题 + 自适应复习 + 机构管理 SaaS"),
        ("明天", "学科知识图谱平台，连接机构、学生、内容"),
        ("未来", "教育领域的 AI 基础设施，重新定义「教」与「学」"),
    ]):
        t = Inches(3.2 + i * 1.0)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.0), t + Inches(0.1), Inches(0.15), Inches(0.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['cyan']; dot.line.fill.background()
        if i < 2:
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.06), t + Inches(0.3), Pt(2), Inches(0.85))
            line.fill.solid(); line.fill.fore_color.rgb = C['border']; line.line.fill.background()
        txt(s, Inches(3.5), t, Inches(1.5), Inches(0.4), stage, size=20, color=C['cyan'], bold=True)
        txt(s, Inches(5.2), t, Inches(6), Inches(0.4), desc, size=17, color=C['muted'])
    txt(s, Inches(2), Inches(6.5), Inches(9.3), Inches(0.3),
        "北京融知高科 · UniMind.ai · korsonedu@gmail.com", size=13, color=C['dim'], align=PP_ALIGN.CENTER)
    page_num(s, 16, T)

    out = os.path.join(os.path.dirname(__file__), 'UniMind_投资版.pptx')
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == '__main__':
    gen_customer()
    gen_investor()
    print("Done!")
